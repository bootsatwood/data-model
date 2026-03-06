"""
Build Internal & External Comms V2 deck
- Uses GTM deck as template (Eventus branding, 14.2x8.0)
- Migrates + enhances comms content
- Adds placeholder slides for Budget Analysis and Drip Campaign
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

GTM_PATH = r"C:\Users\ratwood\OneDrive - Eventus WholeHealth\Vault\05_Board_and_Leadership\GTM Strategy\GTM_Strategy_2026-01-14_V2.pptx"
OUTPUT = r"C:\Users\ratwood\OneDrive - Eventus WholeHealth\Vault\07_Marketing_Enablement\2026-03-02_Internal_External_Comms_V2.pptx"

prs = Presentation(GTM_PATH)

layouts = {}
for i, layout in enumerate(prs.slide_layouts):
    layouts[layout.name] = i
print("Available layouts:", list(layouts.keys()))

# Delete all existing slides
while len(prs.slides._sldIdLst) > 0:
    sldId = prs.slides._sldIdLst[0]
    prs.slides._sldIdLst.remove(sldId)
print(f"Slides after clearing: {len(prs.slides._sldIdLst)}")

BLUE = (0, 70, 127)
GRAY = (120, 120, 120)
LTGRAY = (150, 150, 150)

def e(val):
    return Emu(int(val * 914400))

def add_slide(layout_name):
    idx = layouts.get(layout_name, layouts.get("Content_1_Standard", 11))
    layout = prs.slide_layouts[idx]
    return prs.slides.add_slide(layout)

def set_title(slide, text):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = text
            break

def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(e(left), e(top), e(width), e(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.alignment = align
    if color:
        p.font.color.rgb = RGBColor(*color)
    return txBox

def add_bullets(slide, left, top, width, height, items, size=12, color=None):
    txBox = slide.shapes.add_textbox(e(left), e(top), e(width), e(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.space_after = Pt(4)
        if color:
            p.font.color.rgb = RGBColor(*color)
        pPr = p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, "{%s}buChar" % ns)
        buChar.set("char", "\u2022")
    return txBox

# === SLIDE 1: COVER ===
print("Slide 1: Cover")
s = add_slide("4_Title Slide" if "4_Title Slide" in layouts else "Title Slide")
for ph in s.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "Internal & External\nCommunications"
        for para in ph.text_frame.paragraphs:
            para.font.size = Pt(24)
add_textbox(s, 4.4, 4.9, 8.9, 0.7,
    "Aligning Internal & External Messaging to Drive $7M Annualized Revenue per Quarter",
    size=17, color=(80, 80, 80))
add_textbox(s, 0.2, 4.4, 3.8, 1.2,
    "Powered by MUO Penetration, Metro Alignment, Integrated Services Growth",
    size=14, color=(100, 100, 100))

# === SLIDE 2: AGENDA ===
print("Slide 2: Agenda")
s = add_slide("Content_1_Standard")
set_title(s, "Agenda")
add_bullets(s, 0.8, 1.8, 12.0, 5.5, [
    "2026 Goals & Revenue Targets",
    "Campaign & Communications Overview",
    "Monthly Campaign Calendar",
    "Dual Track: MUO Enterprise + Field/Metro",
    "Content Governance & Workflow",
    "Drip Campaign Strategy",
    "Budget Analysis",
    "Campaign Performance Model",
], size=16)

# === SLIDE 3: 2026 GOALS ===
print("Slide 3: 2026 Goals")
s = add_slide("Content_1_Standard")
set_title(s, "2026 Goals")
add_textbox(s, 0.5, 1.6, 4.0, 0.5, "Business Development", size=16, bold=True, color=BLUE)
add_bullets(s, 0.5, 2.2, 4.0, 2.5,
    ["Support MUO Penetration", "Metro Market Growth", "Increase Lead Generation"], size=13)
add_textbox(s, 5.0, 1.6, 4.0, 0.5, "Learning & Development", size=16, bold=True, color=BLUE)
add_bullets(s, 5.0, 2.2, 4.0, 3.0,
    ["Metro field alignment", "Integrated services expansion (Primary + Behavioral)",
     "Cross-sell primary into mental health accounts", "SS & CCM revenue growth"], size=13)
add_textbox(s, 9.5, 1.6, 4.2, 0.5, "Revenue Targets", size=16, bold=True, color=BLUE)
add_bullets(s, 9.5, 2.2, 4.2, 2.5,
    ["$7M annualized revenue per quarter", "10% marketing-attributed revenue ($700K/qtr)",
     ">=16 MUO touchpoints/month"], size=13)

# === SLIDE 4: CAMPAIGN OVERVIEW ===
print("Slide 4: Campaign Overview")
s = add_slide("Content_1_Standard")
set_title(s, "2026 Campaign & Communications")
add_textbox(s, 0.5, 1.5, 12.5, 0.6,
    "MUO Penetration  |  Metro Alignment  |  Integrated Services Growth",
    size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_textbox(s, 0.5, 2.5, 12.5, 0.5,
    "One unified workflow connecting:", size=14, bold=True, align=PP_ALIGN.CENTER)
add_textbox(s, 0.5, 3.2, 12.5, 0.5,
    "Strategy > Campaign > Content > Enablement > Activation > Revenue",
    size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# === SLIDE 5: CAMPAIGN CALENDAR ===
print("Slide 5: Campaign Calendar")
s = add_slide("Content_1_Standard")
set_title(s, "Monthly Campaign Content & Calendar")
rows, cols = 13, 4
tbl_shape = s.shapes.add_table(rows, cols, e(0.4), e(1.5), e(13.2), e(5.8))
tbl = tbl_shape.table
for j, h in enumerate(["Month", "Theme", "Sales / Revenue Lever", "L&D"]):
    cell = tbl.cell(0, j)
    cell.text = h
    for para in cell.text_frame.paragraphs:
        para.font.size = Pt(11)
        para.font.bold = True
calendar = [
    ("Feb", "Value-Based Care", "Enterprise positioning", ""),
    ("Mar", "Pharmacy + Antipsychotics", "Cross-sell & clinical differentiation", ""),
    ("Apr", "RTH (Return to Hospital)", "Operational & financial value", ""),
    ("May", "Antianxieties", "Behavioral integration growth", ""),
    ("Jun", "Transitions to Care", "Primary penetration & CCM", ""),
    ("Jul", "TBD", "TBD", ""), ("Aug", "TBD", "TBD", ""),
    ("Sep", "TBD", "TBD", ""), ("Oct", "TBD", "TBD", ""),
    ("Nov", "TBD", "TBD", ""), ("Dec", "TBD", "TBD", ""),
    ("Ongoing", "Thought Leadership Series", "LinkedIn/exec positioning", "Workforce, Integration, Financial Future"),
]
for i, (month, theme, lever, ld) in enumerate(calendar, 1):
    for j, val in enumerate([month, theme, lever, ld]):
        cell = tbl.cell(i, j)
        cell.text = val
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(10)

# === SLIDE 6: DUAL TRACK ===
print("Slide 6: Dual Track")
s = add_slide("Content_1_Standard")
set_title(s, "Dual Track: Sales + L&D Alignment")
add_textbox(s, 0.5, 1.5, 6.0, 0.5, "MUO Track (Enterprise)", size=16, bold=True, color=BLUE)
add_textbox(s, 0.5, 2.0, 6.0, 0.4, "Business Development - External Comms", size=11, color=GRAY)
add_bullets(s, 0.5, 2.5, 6.0, 3.0,
    ["White Paper", "Executive Blog", "Sales Executive Summary",
     "LinkedIn Executive Posts", "HubSpot Executive Email"], size=13)
add_textbox(s, 0.5, 5.0, 6.0, 0.5,
    "Purpose: Drive executive-level strategy conversations.", size=12, bold=True)
add_textbox(s, 7.5, 1.5, 6.0, 0.5, "Field / Metro Track", size=16, bold=True, color=BLUE)
add_textbox(s, 7.5, 2.0, 6.0, 0.4, "L&D - Internal Comms", size=11, color=GRAY)
add_bullets(s, 7.5, 2.5, 6.0, 3.0,
    ["2-3 Blogs", "Case Highlight", "Sales One-Pager",
     "Local Email Nurture", "Social Posts"], size=13)
add_textbox(s, 7.5, 5.0, 6.0, 0.5,
    "Purpose: Support BD and cross-sell in corridor markets.", size=12, bold=True)

# === SLIDE 7: GOVERNANCE & WORKFLOW ===
print("Slide 7: Governance & Workflow")
s = add_slide("Content_1_Standard")
set_title(s, "Content Governance & Workflow Model")
add_textbox(s, 0.5, 1.5, 6.5, 0.4, "Quarterly Workflow", size=16, bold=True, color=BLUE)
add_bullets(s, 0.5, 2.0, 6.5, 4.5, [
    "1. Quarterly Campaign Strategy",
    "2. Monthly Campaign Brief",
    "3. Campaign Messaging Guide",
    "4. Content Production (MUO + Field Tracks)",
    "5. Marketing Finalization (SEO/AEO + HubSpot)",
    "6. Internal Enablement Rollout",
    "7. External Activation",
    "8. Revenue + Performance Review",
], size=13)
add_textbox(s, 7.5, 1.5, 6.0, 0.4, "Key Distinctions", size=16, bold=True, color=BLUE)
add_bullets(s, 7.5, 2.0, 6.0, 2.0, [
    "Comms = Authorship", "Marketing = Market Readiness",
    "Strategy = Governance", "Sales = Activation",
], size=13)
add_textbox(s, 7.5, 4.0, 6.0, 0.4, "Pre-Launch Enablement Kit", size=14, bold=True, color=BLUE)
add_bullets(s, 7.5, 4.5, 6.0, 2.5, [
    "1-page Messaging Guide", "Approved proof points",
    "Elevator pitch + FAQ sheet", "Sales slide + Clinician summary",
], size=12)
add_textbox(s, 7.5, 6.5, 6.0, 0.5,
    "Goal: Internal language mirrors external language.", size=12, bold=True, color=BLUE)

# === SLIDE 8: DRIP CAMPAIGN (PLACEHOLDER) ===
print("Slide 8: Drip Campaign (placeholder)")
s = add_slide("Content_1_Standard")
set_title(s, "Drip Campaign Strategy")
add_textbox(s, 0.5, 1.5, 12.5, 0.5,
    "Definition & Cadence by Channel", size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_textbox(s, 0.5, 2.5, 12.5, 1.0,
    "[PLACEHOLDER] Define the automated, sequenced communication strategy that nurtures "
    "prospects through the sales funnel. Reference Michelle's cadence vision for frequency by channel.",
    size=14, color=LTGRAY)
add_textbox(s, 0.5, 3.8, 6.0, 0.4, "Content to Develop:", size=14, bold=True)
add_bullets(s, 0.5, 4.3, 6.0, 3.5, [
    "Drip campaign definition and objectives",
    "Channel cadence (email, LinkedIn, blog, events)",
    "Sequence triggers and timing",
    "Content mapping to buyer journey stage",
    "HubSpot automation workflow",
    "Metrics: open rates, CTR, conversion targets",
], size=12)
add_textbox(s, 7.5, 3.8, 6.0, 0.4, "Reference Inputs:", size=14, bold=True)
add_bullets(s, 7.5, 4.3, 6.0, 2.5, [
    "Michelle Cadence Email (Archive)",
    "Michelle 2026 Marketing Plan (Archive)",
    "Campaigns Example Focus - Brooke (Archive)",
    "HubSpot workflow data mapping (08_GTM_Systems)",
], size=12)

# === SLIDE 9: BUDGET ANALYSIS (PLACEHOLDER) ===
print("Slide 9: Budget Analysis (placeholder)")
s = add_slide("Content_1_Standard")
set_title(s, "Budget Analysis")
add_textbox(s, 0.5, 1.5, 12.5, 0.5,
    "Marketing & BD Spend Allocation", size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_textbox(s, 0.5, 2.5, 12.5, 1.0,
    "[PLACEHOLDER] Summarize BD spend analysis findings and marketing budget allocation "
    "for 2026 campaigns. Reference BD_Spend_Analysis folder for transaction logs and observations.",
    size=14, color=LTGRAY)
add_textbox(s, 0.5, 3.8, 6.0, 0.4, "Content to Develop:", size=14, bold=True)
add_bullets(s, 0.5, 4.3, 6.0, 3.0, [
    "2025 BD spend summary and ROI observations",
    "2026 marketing budget allocation by channel",
    "Campaign-level budget breakdown",
    "Cost per lead / cost per acquisition targets",
    "Investment vs expected revenue impact",
], size=12)
add_textbox(s, 7.5, 3.8, 6.0, 0.4, "Reference Inputs:", size=14, bold=True)
add_bullets(s, 7.5, 4.3, 6.0, 2.0, [
    "2025 BD Reconciled Transaction Log",
    "2025 BD Spend Observations & Questions",
    "2025 BD Transaction Log",
], size=12)

# === SLIDE 10: PERFORMANCE MODEL ===
print("Slide 10: Performance Model")
s = add_slide("Content_1_Standard")
set_title(s, "Campaign Performance Model")
add_textbox(s, 0.5, 1.5, 4.0, 0.5, "Revenue Metrics (Sales)", size=14, bold=True, color=BLUE)
add_bullets(s, 0.5, 2.1, 4.0, 2.5, [
    "$7M annualized revenue target", "10% marketing-sourced revenue",
    "Primary expansion rate", "Cross-sell conversion",
], size=12)
add_textbox(s, 5.0, 1.5, 4.0, 0.5, "Engagement Metrics (Mktg)", size=14, bold=True, color=BLUE)
add_bullets(s, 5.0, 2.1, 4.0, 2.5, [
    "White paper downloads", "Executive engagement",
    "HubSpot lead generation", "Campaign-referenced meetings",
], size=12)
add_textbox(s, 9.5, 1.5, 4.0, 0.5, "Behavior Metrics (L&D)", size=14, bold=True, color=BLUE)
add_bullets(s, 9.5, 2.1, 4.0, 2.5, [
    "Campaign rollout completion", "Message adoption in field",
    "Intranet engagement", "Internal email open rates", "Usage of campaign slides",
], size=12)
add_textbox(s, 0.5, 5.0, 4.0, 0.4, "BD Dashboard", size=12, bold=True, align=PP_ALIGN.CENTER)
add_textbox(s, 5.0, 5.0, 4.0, 0.4, "Marketing Dashboard", size=12, bold=True, align=PP_ALIGN.CENTER)
add_textbox(s, 9.5, 5.0, 4.0, 0.4, "L&D Dashboard", size=12, bold=True, align=PP_ALIGN.CENTER)
add_textbox(s, 0.5, 5.5, 4.0, 0.6, ">=16 MUO touchpoints/month\nMetro penetration tracking", size=11)
add_textbox(s, 5.0, 5.5, 4.0, 0.6, "Leads generated / Downloads\nEmail engagement / Conversion", size=11)
add_textbox(s, 9.5, 5.5, 4.0, 0.6, "Enablement attendance\nMessage recall / Asset usage %", size=11)

# === SAVE ===
prs.save(OUTPUT)
print(f"\nSaved: {OUTPUT}")
print(f"Total slides: {len(list(prs.slides))}")
